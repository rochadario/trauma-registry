"""
create_pitch_en.py
English version of RESPOND Guatemala pitch deck — larger fonts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = '/Users/dariorocha/Claude/trauma-registry/RESPOND_Guatemala_Platform_Pitch_EN.pptx'

NAVY        = RGBColor(0x1B, 0x3A, 0x6B)
BLUE        = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT_BLUE  = RGBColor(0xEB, 0xF5, 0xFB)
MED_BLUE    = RGBColor(0xD6, 0xEA, 0xF8)
GREEN       = RGBColor(0x1E, 0x8B, 0x4C)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
RED         = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED   = RGBColor(0xFD, 0xED, 0xEC)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_ORANGE= RGBColor(0xFE, 0xF9, 0xE7)
YELLOW      = RGBColor(0xF3, 0x9C, 0x12)
DARK        = RGBColor(0x2C, 0x3E, 0x50)
GRAY        = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY  = RGBColor(0xF2, 0xF3, 0xF4)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEAL        = RGBColor(0x0E, 0x86, 0x8A)
LIGHT_TEAL  = RGBColor(0xD1, 0xF2, 0xEB)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def add_slide(): return prs.slides.add_slide(blank)

def rect(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); return s

def rect_border(sl, x, y, w, h, fill, bc, bp=1.5):
    s = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = bc; s.line.width = Pt(bp); return s

def txt(sl, text, x, y, w, h, size=14, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return txb

def txt_ml(sl, text, x, y, w, h, size=14, bold=False, color=DARK,
           align=PP_ALIGN.LEFT, italic=False):
    txb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color; r.font.italic = italic
    return txb

def bg(sl, color): rect(sl, 0, 0, 13.33, 7.5, color)

def header(sl, title, sub=None):
    rect(sl, 0, 0, 13.33, 1.3, NAVY)
    txt(sl, title, 0.3, 0.1, 12.5, 0.7, size=26, bold=True, color=WHITE)
    if sub:
        txt(sl, sub, 0.3, 0.82, 12.5, 0.44, size=14, color=RGBColor(0xAE,0xD6,0xF1), italic=True)

def footer(sl, t='RESPOND Guatemala · Trauma Registry · Confidential'):
    rect(sl, 0, 7.15, 13.33, 0.35, LIGHT_GRAY)
    txt(sl, t, 0.3, 7.17, 12.7, 0.28, size=10, color=GRAY)

def cost_badge(sl, amount, label, x, y, color=GREEN):
    rect(sl, x, y, 2.5, 1.0, color)
    txt(sl, amount, x, y+0.04, 2.5, 0.6, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, label, x, y+0.63, 2.5, 0.34, size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NAVY)
rect(sl, 0, 5.85, 13.33, 1.65, RGBColor(0x12,0x27,0x4A))
rect(sl, 0, 0, 0.18, 7.5, BLUE)
rect(sl, 0, 5.8, 13.33, 0.08, BLUE)

txt(sl, 'RESPOND Guatemala', 0.5, 0.5, 12.0, 0.9, size=46, bold=True, color=WHITE)
txt(sl, 'Digital Trauma Registry', 0.5, 1.45, 12.0, 0.7, size=30, color=RGBColor(0xAE,0xD6,0xF1))
txt(sl, '─────────────────────────────────────────────────────',
    0.5, 2.18, 12.0, 0.4, size=14, color=RGBColor(0x5D,0x8A,0xC4))
txt(sl, 'Platform Overview  ·  HIPAA Compliance  ·  Backend Options  ·  Scaling Strategy',
    0.5, 2.62, 12.0, 0.52, size=17, color=RGBColor(0xAE,0xD6,0xF1), italic=True)

stats = [('16', 'form steps'), ('2', 'languages EN/ES'), ('$0', 'current cost/mo'), ('∞', 'patients')]
for i, (val, lbl) in enumerate(stats):
    bx = 0.5 + i * 3.2
    rect(sl, bx, 3.45, 2.8, 1.25, RGBColor(0x1E,0x49,0x8A))
    txt(sl, val, bx, 3.5, 2.8, 0.72, size=36, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(sl, lbl, bx, 4.18, 2.8, 0.44, size=12, color=RGBColor(0xAE,0xD6,0xF1), align=PP_ALIGN.CENTER)

txt(sl, 'trauma-registry.vercel.app/en/patients', 0.5, 6.0, 8.0, 0.48, size=14,
    color=RGBColor(0x5D,0x8A,0xC4), italic=True)
txt(sl, 'March 2026', 10.5, 6.0, 2.5, 0.48, size=14, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS IT?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'What is RESPOND Guatemala?',
       'A web platform for prospective trauma registry in resource-limited hospitals')

rect(sl, 0.3, 1.45, 5.9, 5.7, LIGHT_RED)
txt(sl, 'The Problem', 0.5, 1.58, 5.5, 0.48, size=17, bold=True, color=RED)
problems = [
    'No digital trauma registry systems in Guatemalan hospitals',
    'Paper-based records prevent analysis and quality improvement',
    'Without data there is no evidence to request resources',
    'Commercial solutions cost $50,000–$200,000/year',
    'Existing systems require constant internet connectivity',
]
for i, p in enumerate(problems):
    txt(sl, f'▸  {p}', 0.45, 2.22+i*0.88, 5.55, 0.78, size=13, color=DARK)

rect(sl, 6.5, 1.45, 6.5, 5.7, LIGHT_GREEN)
txt(sl, 'Our Solution', 6.7, 1.58, 6.1, 0.48, size=17, bold=True, color=GREEN)
solutions = [
    'Free web app, accessible from any device with a browser',
    'Works offline (PWA) — syncs when internet returns',
    '16-step form optimized for real clinical workflow',
    'Automatic GCS and ISS calculation — no manual tables',
    'Real-time dashboard for the trauma team',
    'Bilingual English/Spanish — ready for immediate use',
    'Current operating cost: $0/month',
]
for i, s in enumerate(solutions):
    txt(sl, f'✓  {s}', 6.55, 2.22+i*0.72, 6.3, 0.65, size=13, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT DOES IT CONTAIN?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Platform Features', 'Current MVP — all modules live at trauma-registry.vercel.app')

features = [
    ('📋', '16-Step\nClinical Form', LIGHT_BLUE),
    ('🧠', 'Automatic\nGCS Calculator', LIGHT_GREEN),
    ('📊', 'Auto ISS\nScore', LIGHT_BLUE),
    ('🗺️', 'Interactive\nBody Map', LIGHT_GREEN),
    ('📈', 'Real-Time\nDashboard', LIGHT_BLUE),
    ('📥', 'Excel\nExport', LIGHT_GREEN),
    ('📴', 'Offline Mode\n(PWA)', LIGHT_BLUE),
    ('🌐', 'Bilingual\nEN / ES', LIGHT_GREEN),
    ('🔐', 'Role-Based\nAuth', LIGHT_BLUE),
    ('🏥', 'Multi-Hospital\nReady', LIGHT_GREEN),
    ('📱', 'Mobile\nFriendly', LIGHT_BLUE),
    ('🔄', 'Auto\nSync', LIGHT_GREEN),
]
cols = 6
for i, (emoji, label, color) in enumerate(features):
    col = i % cols; row = i // cols
    x = 0.3 + col * 2.16; y = 1.5 + row * 2.1
    rect_border(sl, x, y, 1.98, 1.88, color, BLUE, 1)
    txt(sl, emoji, x, y+0.1, 1.98, 0.72, size=30, align=PP_ALIGN.CENTER)
    txt_ml(sl, label, x, y+0.95, 1.98, 0.85, size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW WAS IT BUILT
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Technical Architecture',
       'Modern production-grade stack — $0/month operating cost')

rect(sl, 0.3, 1.45, 5.9, 5.7, LIGHT_GRAY)
txt(sl, 'Tech Stack', 0.5, 1.58, 5.5, 0.48, size=17, bold=True, color=NAVY)

stack = [
    ('⚡ Next.js 14', 'Full-stack React framework (frontend + backend)'),
    ('🔷 TypeScript', 'Static typing — fewer bugs in production'),
    ('🗄️ Supabase', 'PostgreSQL database + authentication'),
    ('🎨 Tailwind CSS', 'Responsive mobile-first design'),
    ('🧩 Shadcn/UI', 'Accessible, consistent UI components'),
    ('☁️ Vercel', 'Global hosting, CDN, auto-deploys'),
    ('📴 Service Worker', 'Offline-first PWA with background sync'),
]
for i, (tech, desc) in enumerate(stack):
    rect(sl, 0.35, 2.18+i*0.58, 5.8, 0.53, WHITE)
    txt(sl, tech, 0.45, 2.21+i*0.58, 2.2, 0.46, size=13, bold=True, color=NAVY)
    txt(sl, desc, 2.72, 2.23+i*0.58, 3.35, 0.42, size=12, color=DARK)

rect(sl, 6.5, 1.45, 6.5, 5.7, LIGHT_BLUE)
txt(sl, 'Current Data Flow', 6.7, 1.58, 6.1, 0.48, size=17, bold=True, color=NAVY)

flow = [
    ('👤  User (browser)', NAVY, WHITE),
    ('⬇', None, None),
    ('📱  Next.js App (Vercel)', BLUE, WHITE),
    ('⬇', None, None),
    ('🔐  Supabase Auth  (JWT)', RGBColor(0x1A,0x5A,0x9E), WHITE),
    ('⬇', None, None),
    ('🗄️  Supabase PostgreSQL  (RLS)', RGBColor(0x0E,0x6B,0x5E), WHITE),
    ('⬇', None, None),
    ('📊  Dashboard / Excel Export', NAVY, WHITE),
]
y_pos = 2.15
for item, fill, tc in flow:
    if item == '⬇':
        txt(sl, '↓', 9.25, y_pos, 1.0, 0.35, size=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        y_pos += 0.35
    else:
        rect(sl, 6.65, y_pos, 6.1, 0.65, fill)
        txt(sl, item, 6.72, y_pos+0.1, 5.96, 0.46, size=13, bold=True, color=tc, align=PP_ALIGN.CENTER)
        y_pos += 0.75

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SECURITY
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Platform Security', 'What is already in place vs. what is needed for HIPAA')

rect(sl, 0.3, 1.45, 6.1, 5.7, LIGHT_GREEN)
txt(sl, '✅  Already Implemented', 0.5, 1.58, 5.7, 0.48, size=17, bold=True, color=GREEN)

have = [
    ('🔒 HTTPS/TLS', 'All traffic encrypted in transit — auto cert by Vercel'),
    ('🔐 JWT Auth', 'Signed tokens with automatic rotation (Supabase Auth)'),
    ('🛡️ Row Level Security', 'Each user sees only their own patients (PostgreSQL RLS)'),
    ('💾 Encryption at Rest', 'AES-256 database encryption (Supabase/AWS)'),
    ('🌐 DDoS Protection', 'Vercel Edge Network — included by default'),
    ('🔑 Env Variables', 'Secrets in Vercel, never in source code'),
    ('📦 SOC 2 Type II', 'Supabase certified — annual security audits'),
]
for i, (title, desc) in enumerate(have):
    rect(sl, 0.35, 2.2+i*0.68, 6.0, 0.62, WHITE)
    txt(sl, title, 0.45, 2.23+i*0.68, 2.1, 0.52, size=12, bold=True, color=GREEN)
    txt(sl, desc, 2.62, 2.26+i*0.68, 3.6, 0.45, size=12, color=DARK)

rect(sl, 6.7, 1.45, 6.3, 5.7, LIGHT_ORANGE)
txt(sl, '⚠️  Required for HIPAA', 6.9, 1.58, 5.9, 0.48, size=17, bold=True, color=ORANGE)

need = [
    ('⏱️ Session Timeout', 'Auto-logout after 15 min of inactivity'),
    ('📋 Audit Log', 'Who accessed/modified which patient and when'),
    ('📱 MFA', '2-factor authentication for all users'),
    ('🔏 Signed BAA', 'Legal contract with the database provider'),
    ('📄 Privacy Policy', 'Published Notice of Privacy Practices'),
    ('📊 Risk Assessment', 'Formal security risk assessment document'),
    ('🎓 Training Records', 'Documentation of staff HIPAA training'),
]
for i, (title, desc) in enumerate(need):
    rect(sl, 6.75, 2.2+i*0.68, 6.2, 0.62, WHITE)
    txt(sl, title, 6.85, 2.23+i*0.68, 2.1, 0.52, size=12, bold=True, color=ORANGE)
    txt(sl, desc, 9.0, 2.26+i*0.68, 3.8, 0.45, size=12, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — WHY HIPAA?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Why Do We Need HIPAA Compliance?',
       'Required by mentors — standard for US-affiliated clinical research')

rect(sl, 0.3, 1.45, 4.1, 5.7, LIGHT_RED)
txt(sl, '⚖️  What is HIPAA?', 0.5, 1.58, 3.7, 0.48, size=15, bold=True, color=RED)
txt_ml(sl,
    'Health Insurance Portability\nand Accountability Act\n(US federal law, 1996)\n\n'
    'Protects identifiable patient\nmedical information\n(PHI — Protected Health\nInformation)\n\n'
    'Requires:\n• Confidentiality\n• Integrity\n• Availability\nof all patient data',
    0.45, 2.18, 3.72, 4.7, size=13, color=DARK)

rect(sl, 4.65, 1.45, 4.0, 5.7, LIGHT_ORANGE)
txt(sl, '🌎  Does it Apply\n     in Guatemala?', 4.85, 1.58, 3.7, 0.78, size=15, bold=True, color=ORANGE)
txt_ml(sl,
    'Technically NO —\nHIPAA is US law.\n\n'
    'BUT it applies when:\n\n'
    '• Mentors are at a US\n  institution (BU)\n\n'
    '• Funding is from NIH\n  or US sources\n\n'
    '• Data will be published\n  in journals requiring\n  international standards\n\n'
    '→ HIPAA compliance is\n  the gold standard for\n  serious clinical research',
    4.85, 2.5, 3.72, 4.4, size=13, color=DARK)

rect(sl, 8.9, 1.45, 4.1, 5.7, LIGHT_BLUE)
txt(sl, '📋  What PHI Does\n      It Protect?', 9.1, 1.58, 3.7, 0.78, size=15, bold=True, color=NAVY)
phi = [
    'Patient name',
    'Date of birth',
    'Medical record number',
    'Admission / discharge date',
    'Diagnosis',
    'Clinical results (GCS, ISS)',
    'Procedures performed',
    'Discharge condition',
    'Any data that can\nidentify the patient',
]
for i, item in enumerate(phi):
    txt_ml(sl, f'▸  {item}', 9.1, 2.5+i*0.52, 3.72, 0.48, size=13, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OPTION 1: REDCAP (RECOMMENDED)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Option 1: REDCap Integration  ⭐ RECOMMENDED',
       'The app on Vercel sends data directly to the institutional REDCap server')

rect(sl, 10.7, 1.4, 2.3, 0.58, GREEN)
txt(sl, '⭐ RECOMMENDED', 10.72, 1.43, 2.26, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rect(sl, 0.3, 1.5, 12.73, 2.2, LIGHT_GRAY)
txt(sl, 'Data Flow with REDCap:', 0.5, 1.57, 5.0, 0.4, size=14, bold=True, color=NAVY)

boxes = [
    ('👤\nNurse / MD', NAVY, WHITE, 0.4, 2.0, 1.8, 1.5),
    ('📱\nRESPOND App\n(Vercel)', BLUE, WHITE, 2.55, 2.0, 1.9, 1.5),
    ('🔀\nNext.js\nAPI Route', RGBColor(0x5D,0x6D,0x7E), WHITE, 4.75, 2.0, 1.8, 1.5),
    ('🏛️\nREDCap API\n(institutional)', TEAL, WHITE, 6.85, 2.0, 2.0, 1.5),
    ('🗄️\nREDCap Server\nBU / Institution', RGBColor(0x0E,0x5E,0x5E), WHITE, 9.15, 2.0, 2.1, 1.5),
    ('🔒\nHIPAA\nCovered', GREEN, WHITE, 11.55, 2.0, 1.55, 1.5),
]
for lbl, fill, tc, x, y, w, h in boxes:
    rect(sl, x, y, w, h, fill)
    txt_ml(sl, lbl, x, y+0.12, w, h-0.18, size=12, bold=True, color=tc, align=PP_ALIGN.CENTER)
for ax in [2.3, 4.5, 6.7, 8.95, 11.35]:
    txt(sl, '→', ax, 2.58, 0.35, 0.45, size=20, bold=True, color=BLUE)

rect(sl, 0.3, 3.9, 6.1, 3.25, LIGHT_GREEN)
txt(sl, '✅  Advantages', 0.5, 4.02, 5.7, 0.45, size=16, bold=True, color=GREEN)
pros = [
    ('💰 Cost: $0/year', 'If institution already has REDCap (BU does)'),
    ('🔒 HIPAA covered', 'Institution is legally responsible, not you'),
    ('📊 Native exports', 'SPSS, Stata, R, Excel, CSV — built into REDCap'),
    ('🎓 Research standard', '6,000+ institutions worldwide use REDCap'),
    ('🚀 No BAA to sign', 'Institution IT department manages it'),
]
for i, (t, d) in enumerate(pros):
    txt(sl, t, 0.45, 4.6+i*0.52, 2.1, 0.45, size=13, bold=True, color=GREEN)
    txt(sl, d, 2.62, 4.62+i*0.52, 3.7, 0.42, size=12.5, color=DARK)

rect(sl, 6.7, 3.9, 6.3, 3.25, LIGHT_ORANGE)
txt(sl, '⚠️  Considerations', 6.9, 4.02, 5.9, 0.45, size=16, bold=True, color=ORANGE)
cons = [
    ('🏛️ Needs institution', 'BU or partner institution must open the project'),
    ('🗺️ Field mapping', 'Form fields must map to REDCap variable names'),
    ('📴 Offline limited', 'REDCap API needs connection (solvable with local cache)'),
    ('📈 Basic dashboard', 'REDCap reports are basic — app can pull via API'),
    ('⏱️ Initial setup', '~1 week to coordinate with IT + map variables'),
]
for i, (t, d) in enumerate(cons):
    txt(sl, t, 6.9, 4.6+i*0.52, 2.2, 0.45, size=13, bold=True, color=ORANGE)
    txt(sl, d, 9.15, 4.62+i*0.52, 3.75, 0.42, size=12.5, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — REDCAP TECHNICAL SETUP
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'REDCap — Technical Integration',
       'Minimal code changes to connect the existing app to the REDCap API')

rect(sl, 0.3, 1.45, 5.9, 5.7, LIGHT_BLUE)
txt(sl, '🔧 What Changes in the Code?', 0.5, 1.58, 5.5, 0.48, size=15, bold=True, color=NAVY)

changes = [
    ('No change', 'The 16-step form stays exactly as-is'),
    ('No change', 'Dashboard can keep reading from REDCap API'),
    ('New', 'Next.js API route: receives form → POST to REDCap'),
    ('New', 'Field mapping: app_field_name → redcap_variable'),
    ('New', 'Env vars: REDCAP_API_TOKEN + REDCAP_URL'),
    ('Optional', 'Local cache (IndexedDB) for offline → sync on reconnect'),
]
labels_c = {'No change': GREEN, 'New': BLUE, 'Optional': ORANGE}
for i, (label, desc) in enumerate(changes):
    lc = labels_c.get(label, GRAY)
    rect(sl, 0.35, 2.25+i*0.8, 1.4, 0.62, lc)
    txt(sl, label, 0.37, 2.28+i*0.8, 1.36, 0.55, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, 1.88, 2.32+i*0.8, 4.2, 0.52, size=13, color=DARK)

rect(sl, 6.5, 1.45, 6.5, 5.7, LIGHT_GRAY)
txt(sl, '📋 Setup Process', 6.7, 1.58, 6.1, 0.48, size=15, bold=True, color=NAVY)

steps = [
    ('1', 'Request REDCap project at BU (IT dept.)', '~1–3 days'),
    ('2', 'Define REDCap variables = form fields', '~1 day'),
    ('3', 'Generate API token in REDCap dashboard', '5 min'),
    ('4', 'Write the API route in Next.js', '~4–8 hrs'),
    ('5', 'Map 50–60 form fields to REDCap API', '~4–8 hrs'),
    ('6', 'End-to-end testing (form → REDCap)', '~1 day'),
    ('7', 'Deploy + train clinical team', '~1 day'),
]
for i, (num, step, time) in enumerate(steps):
    rect(sl, 6.55, 2.25+i*0.7, 0.5, 0.58, NAVY)
    txt(sl, num, 6.57, 2.28+i*0.7, 0.46, 0.5, size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, step, 7.16, 2.32+i*0.7, 4.1, 0.46, size=13, color=DARK)
    txt(sl, time, 11.3, 2.32+i*0.7, 1.55, 0.46, size=12, color=GRAY, italic=True, align=PP_ALIGN.RIGHT)

rect(sl, 6.55, 7.1-0.55, 6.4, 0.5, GREEN)
txt(sl, '⏱️  Total estimated time: 1–2 weeks  ·  Minimal development cost',
    6.65, 7.13-0.55, 6.2, 0.4, size=12.5, bold=True, color=WHITE)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OPTION 2: SUPABASE PRO
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Option 2: Supabase Pro + BAA',
       'Upgrade the existing backend — zero code changes, HIPAA compliant')

cost_badge(sl, '$300', 'per year', 0.3, 1.45, BLUE)
rect(sl, 3.05, 1.45, 10.0, 1.0, LIGHT_BLUE)
txt_ml(sl,
    'The app is already 100% built on Supabase. Simply upgrade to Pro ($25/mo) and sign the BAA.\n'
    'No code migration required. HIPAA coverage activated in days.',
    3.22, 1.52, 9.7, 0.88, size=14, color=DARK)

rect(sl, 0.3, 2.6, 12.73, 1.55, LIGHT_GRAY)
txt(sl, 'Architecture (unchanged):', 0.5, 2.68, 4.5, 0.4, size=14, bold=True, color=NAVY)
flow2 = [
    ('👤 User', NAVY), ('📱 Vercel App', BLUE), ('🔐 Supabase Auth', TEAL),
    ('🗄️ Supabase PostgreSQL', RGBColor(0x0E,0x5E,0x5E)), ('✅ BAA = HIPAA', GREEN)
]
for i, (label, fill) in enumerate(flow2):
    rect(sl, 0.3+i*2.56, 3.1, 2.42, 0.78, fill)
    txt(sl, label, 0.33+i*2.56, 3.14, 2.36, 0.7, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.8+i*2.56, 3.35, 0.38, 0.38, size=18, bold=True, color=BLUE)

rect(sl, 0.3, 4.1, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Advantages', 0.5, 4.22, 5.7, 0.45, size=16, bold=True, color=GREEN)
pros2 = [
    'Zero code changes in the existing app',
    'BAA available on Pro — signed in the Supabase dashboard',
    'Supabase runs on AWS us-east-1 — same infra as major hospital systems',
    'Auth, DB, RLS, storage — all HIPAA under one provider',
    'Dashboard, offline PWA and exports continue working as-is',
]
for i, p in enumerate(pros2):
    txt(sl, f'✓  {p}', 0.45, 4.8+i*0.46, 5.7, 0.42, size=13, color=DARK)

rect(sl, 6.7, 4.1, 6.3, 3.05, LIGHT_ORANGE)
txt(sl, '⚠️  Considerations', 6.9, 4.22, 5.9, 0.45, size=16, bold=True, color=ORANGE)
cons2 = [
    'Fixed $25/mo cost even at zero patients',
    'Vercel stays free (PHI does not pass through Vercel)',
    'Must add session timeout + audit log + MFA (~1–2 days)',
    'Administrative docs still needed (risk assessment, privacy policy)',
]
for i, c in enumerate(cons2):
    txt(sl, f'▸  {c}', 6.9, 4.8+i*0.62, 5.9, 0.55, size=13, color=DARK)

rect(sl, 0.3, 6.75-0.5, 12.73, 0.45, NAVY)
txt(sl, '  STEPS:   1. Upgrade to Supabase Pro   →   2. Sign BAA in dashboard   →   3. Add session timeout + MFA + audit log   →   4. Document risk assessment',
    0.4, 6.78-0.5, 12.5, 0.38, size=12, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — OPTION 3: AWS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Option 3: Amazon Web Services (AWS)',
       'Enterprise-grade infrastructure — maximum flexibility and control')

cost_badge(sl, '$600–$1,800', 'per year (variable)', 0.3, 1.45, RGBColor(0xFF,0x99,0x00))
rect(sl, 3.05, 1.45, 10.0, 1.0, LIGHT_ORANGE)
txt_ml(sl,
    'AWS signs BAA on any paid plan. Architecture: RDS PostgreSQL + Cognito (auth) + Amplify (hosting).\n'
    'Significant code migration required. Used by Epic, Cerner, and large hospital systems.',
    3.22, 1.52, 9.7, 0.88, size=14, color=DARK)

rect(sl, 0.3, 2.6, 12.73, 1.55, LIGHT_GRAY)
txt(sl, 'AWS Architecture:', 0.5, 2.68, 4.5, 0.4, size=14, bold=True, color=NAVY)
aws = [
    ('👤 User', NAVY), ('📱 AWS Amplify\n(hosting)', ORANGE),
    ('🔐 Cognito\n(auth)', RGBColor(0xC0,0x50,0x00)),
    ('🗄️ RDS PostgreSQL\n+ KMS encrypt', RGBColor(0x8B,0x35,0x00)),
    ('✅ BAA = HIPAA', GREEN)
]
for i, (lbl, fill) in enumerate(aws):
    rect(sl, 0.3+i*2.56, 3.08, 2.42, 0.85, fill)
    txt_ml(sl, lbl, 0.33+i*2.56, 3.12, 2.36, 0.78, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.8+i*2.56, 3.35, 0.38, 0.38, size=18, bold=True, color=ORANGE)

rect(sl, 0.3, 4.1, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Advantages', 0.5, 4.22, 5.7, 0.45, size=16, bold=True, color=GREEN)
pros3 = [
    'Gold standard for healthcare IT in the United States',
    'BAA available — AWS is HIPAA-eligible across all key services',
    'Maximum infrastructure control and visibility',
    'Native audit logging via CloudTrail',
    'Scales to millions of patients without re-architecting',
]
for i, p in enumerate(pros3):
    txt(sl, f'✓  {p}', 0.45, 4.8+i*0.46, 5.7, 0.42, size=13, color=DARK)

rect(sl, 6.7, 4.1, 6.3, 3.05, LIGHT_RED)
txt(sl, '✗  Disadvantages', 6.9, 4.22, 5.9, 0.45, size=16, bold=True, color=RED)
cons3 = [
    'Full code migration: auth, DB, API — weeks of work',
    'High operational complexity: IAM, VPC, security groups',
    'Requires DevOps expertise to maintain',
    'Variable cost — can increase with usage',
    'Overkill for a single-hospital pilot stage',
]
for i, c in enumerate(cons3):
    txt(sl, f'✗  {c}', 6.9, 4.8+i*0.46, 5.9, 0.42, size=13, color=DARK)

rect(sl, 0.3, 6.75-0.5, 12.73, 0.45, RGBColor(0xCC,0x77,0x00))
txt(sl, '  Cost breakdown: RDS db.t3.micro ~$15/mo · Cognito <1K users free · Amplify ~$1–5/mo · Extras (backups, KMS, CloudWatch) ~$5–10/mo → Total ~$25–50/mo',
    0.4, 6.78-0.5, 12.5, 0.38, size=12, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — OPTION 4: GOOGLE CLOUD
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Option 4: Google Cloud Platform / Firebase',
       'Low cost at small scale — BAA available on Blaze plan')

cost_badge(sl, '$0–$360', 'per year (pay-as-you-go)', 0.3, 1.45, TEAL)
rect(sl, 3.05, 1.45, 10.0, 1.0, LIGHT_TEAL)
txt_ml(sl,
    'Google signs BAA under GCP. Firebase Auth + Firestore (or Cloud SQL) are HIPAA-eligible.\n'
    'Near-zero cost at low volumes. Significant code migration required.',
    3.22, 1.52, 9.7, 0.88, size=14, color=DARK)

rect(sl, 0.3, 2.6, 12.73, 1.55, LIGHT_GRAY)
txt(sl, 'GCP Architecture:', 0.5, 2.68, 4.5, 0.4, size=14, bold=True, color=NAVY)
gcp = [
    ('👤 User', NAVY), ('📱 Cloud Run\n(hosting)', TEAL),
    ('🔐 Firebase Auth\n(Google)', RGBColor(0x0A,0x70,0x70)),
    ('🗄️ Firestore or\nCloud SQL', RGBColor(0x06,0x50,0x50)),
    ('✅ BAA = HIPAA', GREEN)
]
for i, (lbl, fill) in enumerate(gcp):
    rect(sl, 0.3+i*2.56, 3.08, 2.42, 0.85, fill)
    txt_ml(sl, lbl, 0.33+i*2.56, 3.12, 2.36, 0.78, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.8+i*2.56, 3.35, 0.38, 0.38, size=18, bold=True, color=TEAL)

rect(sl, 0.3, 4.1, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Advantages', 0.5, 4.22, 5.7, 0.45, size=16, bold=True, color=GREEN)
pros4 = [
    'Near-zero cost at low data volumes',
    'Google signs BAA on Blaze plan (free or pay-as-you-go)',
    'Firebase Auth is robust with MFA built-in',
    'Cloud SQL = PostgreSQL (easier to migrate than Firestore)',
    'Google Workspace / Drive integration for reports',
]
for i, p in enumerate(pros4):
    txt(sl, f'✓  {p}', 0.45, 4.8+i*0.46, 5.7, 0.42, size=13, color=DARK)

rect(sl, 6.7, 4.1, 6.3, 3.05, LIGHT_RED)
txt(sl, '✗  Disadvantages', 6.9, 4.22, 5.9, 0.45, size=16, bold=True, color=RED)
cons4 = [
    'Full code migration from Supabase to Firebase/GCP',
    'Firestore is NoSQL — very different from current schema',
    'Cloud SQL needs more configuration than Supabase',
    'GCP learning curve vs. simplicity of Supabase',
    'Vendor lock-in to Google ecosystem',
]
for i, c in enumerate(cons4):
    txt(sl, f'✗  {c}', 6.9, 4.8+i*0.46, 5.9, 0.42, size=13, color=DARK)

rect(sl, 0.3, 6.75-0.5, 12.73, 0.45, TEAL)
txt(sl, '  Cost: Firestore <50K reads/day free · Firebase Auth free · Cloud Run $0–5/mo at low volume → Total ~$0–30/mo',
    0.4, 6.78-0.5, 12.5, 0.38, size=12, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COST COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Cost & Feature Comparison',
       'Annual analysis for a single-hospital pilot (~500–1,000 patients/year)')

options = [
    ('REDCap\n(Institutional)', '$0', '/year', GREEN),
    ('Supabase Pro', '$300', '/year', BLUE),
    ('Google Cloud\n(Firebase)', '$0–360', '/year', TEAL),
    ('AWS\n(RDS + Cognito)', '$300–600', '/year', ORANGE),
]
for i, (name, cost, unit, color) in enumerate(options):
    x = 0.3 + i * 3.26
    rect(sl, x, 1.4, 3.05, 1.12, color)
    txt_ml(sl, name, x, 1.42, 3.05, 0.44, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, cost, x, 1.83, 3.05, 0.56, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, unit, x, 2.38, 3.05, 0.3, size=11, color=RGBColor(0xDD,0xDD,0xDD), align=PP_ALIGN.CENTER)

headers_t = ['Feature', 'REDCap', 'Supabase Pro', 'Google Cloud', 'AWS']
col_x = [0.3, 3.57, 6.07, 8.57, 11.07]
col_w = [3.22, 2.45, 2.45, 2.45, 2.2]

rect(sl, 0.3, 2.72, 12.73, 0.5, NAVY)
for j, (h, x, w) in enumerate(zip(headers_t, col_x, col_w)):
    txt(sl, h, x+0.06, 2.75, w-0.1, 0.42, size=13, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

rows = [
    ('HIPAA Compliance',         '✅ Institutional',  '✅ BAA on Pro',    '✅ BAA Blaze',    '✅ BAA any paid plan'),
    ('Annual cost estimate',     '💚 $0',             '🔵 $300',          '🟢 $0–$360',      '🟡 $300–$600'),
    ('Code changes required',    '🟡 Medium (API)',   '✅ Minimal',       '🔴 High',         '🔴 High'),
    ('Time to implement',        '1–2 weeks',         '2–3 days',         '4–8 weeks',       '4–8 weeks'),
    ('Offline / PWA',            '🟡 With cache',     '✅ Full support',  '🟡 With cache',   '🟡 With cache'),
    ('Real-time dashboard',      '🟡 API pull',       '✅ Works now',     '🟡 Needs dev',    '🟡 Needs dev'),
    ('Research data exports',    '✅ SPSS/SAS/R/CSV', '✅ Excel/CSV',     '🟡 CSV/JSON',     '🟡 CSV/JSON'),
    ('Known by researchers',     '✅ Gold standard',  '🟡 Unknown',       '🟡 Unknown',      '🟡 Unknown'),
    ('Scalability',              '✅ High',           '✅ High',          '✅ High',         '✅ Very high'),
]
row_colors = [LIGHT_GRAY, WHITE]
for i, row in enumerate(rows):
    bg_c = row_colors[i % 2]
    rect(sl, 0.3, 3.26+i*0.41, 12.73, 0.4, bg_c)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        txt(sl, cell, x+0.06, 3.28+i*0.41, w-0.1, 0.35, size=11, color=DARK, align=al)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RECOMMENDATION
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Recommendation', 'Two-phase strategy based on project goals')

rect(sl, 0.3, 1.45, 6.1, 5.7, LIGHT_GREEN)
rect(sl, 0.3, 1.45, 6.1, 0.65, GREEN)
txt(sl, '⭐  Phase 1 — Research Pilot', 0.45, 1.49, 5.8, 0.55, size=16, bold=True, color=WHITE)
txt(sl, 'REDCap + RESPOND Guatemala App', 0.45, 2.25, 5.8, 0.5, size=17, bold=True, color=NAVY)
txt(sl, '$0/year if BU has REDCap', 0.45, 2.76, 5.8, 0.4, size=14, color=GREEN, bold=True)

rec1 = [
    '✓  HIPAA covered by institution — zero cost',
    '✓  Data on BU servers — full legal protection',
    '✓  REDCap is the standard in clinical publications',
    '✓  Direct export to SPSS/Stata/R for analysis',
    '✓  BU IRB can approve directly',
    '✓  Minimal monthly operating cost',
    '✓  1–2 weeks technical integration',
]
for i, item in enumerate(rec1):
    txt(sl, item, 0.45, 3.3+i*0.54, 5.7, 0.48, size=13, color=DARK)

rect(sl, 6.7, 1.45, 6.3, 5.7, LIGHT_BLUE)
rect(sl, 6.7, 1.45, 6.3, 0.65, BLUE)
txt(sl, '🚀  Phase 2 — Regional Scale', 6.85, 1.49, 6.0, 0.55, size=16, bold=True, color=WHITE)
txt(sl, 'Supabase Pro (or AWS with funding)', 6.85, 2.25, 6.0, 0.5, size=16, bold=True, color=NAVY)
txt(sl, '$300/year (Supabase) · $300–600/year (AWS)', 6.85, 2.76, 6.0, 0.4, size=13, color=BLUE, bold=True)

rec2 = [
    '✓  When scaling to multiple hospitals (>3)',
    '✓  If advanced real-time dashboard is needed',
    '✓  Once the project secures formal funding',
    '✓  App is already built — just upgrade + sign BAA',
    '✓  Independent from any academic institution',
    '✓  Full ownership of patient data',
    '✓  2–3 days to implement from today',
]
for i, item in enumerate(rec2):
    txt(sl, item, 6.85, 3.3+i*0.54, 6.1, 0.48, size=13, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — NEXT STEPS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
header(sl, 'Next Steps', 'Path to launching the pilot with full HIPAA compliance')

steps_main = [
    (NAVY,  '1', 'Confirm REDCap Access',
     'Ask BU IT if a REDCap project is available for Guatemala.\nAlternative: submit a new project request (form + approval ~1–3 days).'),
    (BLUE,  '2', 'IRB / Ethics Committee',
     'Submit research protocol to BU IRB.\nREDCap integration facilitates approval — it is BU\'s standard platform.'),
    (TEAL,  '3', 'REDCap Technical Integration',
     'Map the 16-step form fields to REDCap variables.\nWrite the Next.js API route to send data (~4–8 hours of development).'),
    (GREEN, '4', 'Add HIPAA Technical Controls',
     'Session timeout (15 min) · Audit log · MFA for all users.\nEstimated: 1–2 days of development, independent of backend choice.'),
    (ORANGE,'5', 'Pilot Hospital + Training',
     'Add 1 hospital to the dropdown · Create 3–5 user accounts.\nClinical team training: ~2 hours.'),
]
for i, (color, num, title, desc) in enumerate(steps_main):
    y = 1.5 + i * 1.12
    rect(sl, 0.3, y, 0.58, 1.0, color)
    txt(sl, num, 0.3, y+0.15, 0.58, 0.72, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, 0.93, y, 12.0, 1.0, LIGHT_GRAY if i % 2 == 0 else WHITE)
    txt(sl, title, 1.05, y+0.06, 5.5, 0.46, size=15, bold=True, color=color)
    txt_ml(sl, desc, 1.05, y+0.54, 11.75, 0.44, size=12.5, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 0.18, 7.5, BLUE)
rect(sl, 0, 5.95, 13.33, 0.08, BLUE)

txt(sl, 'RESPOND Guatemala', 0.5, 0.75, 12.0, 0.85, size=40, bold=True, color=WHITE)
txt(sl, 'Digital Trauma Registry · Ready for Pilot', 0.5, 1.65, 12.0, 0.58, size=22,
    color=RGBColor(0xAE,0xD6,0xF1), italic=True)

rect(sl, 0.5, 2.45, 12.0, 2.1, RGBColor(0x1E,0x49,0x8A))
summary = [
    ('✅', 'Complete MVP live in production — trauma-registry.vercel.app'),
    ('✅', 'Core security in place: HTTPS, JWT, RLS, AES-256 encryption, SOC 2'),
    ('✅', 'HIPAA options identified: REDCap ($0) · Supabase Pro ($300/yr) · AWS · GCP'),
    ('⭐', 'Recommendation: institutional REDCap for pilot → zero cost, maximum credibility'),
    ('🚀', 'Time to first registered patient: 2–4 weeks with IRB approval'),
]
for i, (icon, text) in enumerate(summary):
    txt(sl, icon, 0.7, 2.58+i*0.37, 0.48, 0.33, size=14, bold=True, color=YELLOW)
    txt(sl, text, 1.22, 2.6+i*0.37, 11.0, 0.33, size=13, color=WHITE)

txt(sl, 'Questions?', 0.5, 4.75, 12.0, 0.7, size=30, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txt(sl, 'trauma-registry.vercel.app/en/patients', 0.5, 5.48, 12.0, 0.48,
    size=16, color=RGBColor(0x5D,0x8A,0xC4), italic=True, align=PP_ALIGN.CENTER)

rect(sl, 0, 6.1, 13.33, 1.4, RGBColor(0x12,0x27,0x4A))
txt(sl, 'RESPOND Guatemala  ·  Trauma Registry  ·  Confidential  ·  March 2026',
    0.5, 6.6, 12.33, 0.48, size=12, color=GRAY, align=PP_ALIGN.CENTER)

prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Total slides: {len(prs.slides)}')
