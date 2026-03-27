"""
create_pitch.py
Generates RESPOND_Guatemala_Platform_Pitch.pptx
Professional pitch deck: platform overview + HIPAA options
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from lxml import etree
import copy

OUT = '/Users/dariorocha/Claude/trauma-registry/RESPOND_Guatemala_Platform_Pitch.pptx'

# ── Colors ────────────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x3A, 0x6B)
BLUE       = RGBColor(0x2E, 0x6D, 0xB4)
LIGHT_BLUE = RGBColor(0xEB, 0xF5, 0xFB)
MED_BLUE   = RGBColor(0xD6, 0xEA, 0xF8)
GREEN      = RGBColor(0x1E, 0x8B, 0x4C)
LIGHT_GREEN= RGBColor(0xD5, 0xF5, 0xE3)
RED        = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED  = RGBColor(0xFD, 0xED, 0xEC)
ORANGE     = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_ORANGE=RGBColor(0xFE, 0xF9, 0xE7)
YELLOW     = RGBColor(0xF3, 0x9C, 0x12)
DARK       = RGBColor(0x2C, 0x3E, 0x50)
GRAY       = RGBColor(0x7F, 0x8C, 0x8D)
LIGHT_GRAY = RGBColor(0xF2, 0xF3, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEAL       = RGBColor(0x0E, 0x86, 0x8A)
LIGHT_TEAL = RGBColor(0xD1, 0xF2, 0xEB)
PURPLE     = RGBColor(0x6C, 0x3A, 0x83)
LIGHT_PURPLE=RGBColor(0xF4, 0xEC, 0xF7)

# ── Slide size: 13.33 x 7.5 (widescreen 16:9) ─────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank_layout)

# ── Drawing helpers ────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, radius=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def rect_border(slide, x, y, w, h, fill, border_color, border_pt=1.5):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_pt)
    return shape

def txt(slide, text, x, y, w, h, size=13, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.italic = italic
    return txb

def txt_box(slide, text, x, y, w, h, size=13, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, italic=False):
    """Multi-line text with \n support"""
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.italic = italic
    return txb

def line(slide, x1, y1, x2, y2, color=GRAY, width_pt=1.5):
    from pptx.util import Inches as I
    connector = slide.shapes.add_connector(1, I(x1), I(y1), I(x2), I(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)

def img(slide, path, x, y, w):
    try:
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))
    except:
        pass

def bg(slide, color):
    rect(slide, 0, 0, 13.33, 7.5, color)

def slide_header(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.25, NAVY)
    txt(slide, title, 0.3, 0.12, 12.0, 0.65, size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.78, 12.0, 0.42, size=13, color=RGBColor(0xAE, 0xD6, 0xF1), italic=True)

def footer(slide, text='RESPOND Guatemala · Registro de Trauma · Confidencial'):
    rect(slide, 0, 7.15, 13.33, 0.35, LIGHT_GRAY)
    txt(slide, text, 0.3, 7.17, 12.7, 0.28, size=9, color=GRAY)

def circle(slide, x, y, d, fill):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def icon_box(slide, emoji, label, x, y, bg_color, text_color=DARK, size=12):
    rect_border(slide, x, y, 2.0, 1.45, bg_color, BLUE, 1)
    txt(slide, emoji, x, y+0.05, 2.0, 0.6, size=26, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y+0.75, 2.0, 0.65, size=size, bold=True, color=text_color,
        align=PP_ALIGN.CENTER)

def arrow_right(slide, x, y, color=BLUE):
    """Draw a right-pointing arrow at position"""
    shape = slide.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.4), Inches(0.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def flow_box(slide, text, x, y, w=2.2, h=0.7, fill=LIGHT_BLUE, border=BLUE, tsize=11):
    rect_border(slide, x, y, w, h, fill, border, 1.5)
    txt(slide, text, x+0.05, y+0.08, w-0.1, h-0.15, size=tsize, bold=True,
        color=NAVY, align=PP_ALIGN.CENTER)

def check_row(slide, text, x, y, w=5.5, check_color=GREEN, size=11.5):
    txt(slide, '✓', x, y, 0.3, 0.35, size=size, bold=True, color=check_color)
    txt(slide, text, x+0.3, y, w-0.3, 0.35, size=size, color=DARK)

def x_row(slide, text, x, y, w=5.5, size=11.5):
    txt(slide, '✗', x, y, 0.3, 0.35, size=size, bold=True, color=RED)
    txt(slide, text, x+0.3, y, w-0.3, 0.35, size=size, color=DARK)

def cost_badge(slide, amount, label, x, y, color=GREEN):
    rect(slide, x, y, 2.2, 0.9, color)
    txt(slide, amount, x, y+0.03, 2.2, 0.5, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y+0.55, 2.2, 0.32, size=10, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NAVY)
rect(sl, 0, 5.8, 13.33, 1.7, RGBColor(0x12, 0x27, 0x4A))
rect(sl, 0, 0, 0.18, 7.5, BLUE)
rect(sl, 0, 5.75, 13.33, 0.08, BLUE)

txt(sl, 'RESPOND Guatemala', 0.5, 0.55, 12.0, 0.8, size=42, bold=True, color=WHITE)
txt(sl, 'Registro Digital de Trauma', 0.5, 1.38, 12.0, 0.65, size=28, color=RGBColor(0xAE, 0xD6, 0xF1))
txt(sl, '─────────────────────────────────────────────────────',
    0.5, 2.05, 12.0, 0.4, size=14, color=RGBColor(0x5D, 0x8A, 0xC4))

txt(sl, 'Arquitectura de la plataforma  ·  Cumplimiento HIPAA  ·  Opciones de escalamiento',
    0.5, 2.5, 12.0, 0.5, size=16, color=RGBColor(0xAE, 0xD6, 0xF1), italic=True)

# Stat pills
stats = [('16', 'pasos del formulario'), ('2', 'idiomas ES/EN'), ('0', 'costo actual/mes'), ('∞', 'pacientes')]
for i, (val, lbl) in enumerate(stats):
    bx = 0.5 + i * 3.2
    rect(sl, bx, 3.4, 2.8, 1.2, RGBColor(0x1E, 0x49, 0x8A))
    txt(sl, val, bx, 3.45, 2.8, 0.68, size=34, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    txt(sl, lbl, bx, 4.1, 2.8, 0.42, size=11, color=RGBColor(0xAE, 0xD6, 0xF1), align=PP_ALIGN.CENTER)

txt(sl, 'trauma-registry.vercel.app/en/patients', 0.5, 5.95, 8.0, 0.45, size=13,
    color=RGBColor(0x5D, 0x8A, 0xC4), italic=True)
txt(sl, 'Marzo 2026', 10.5, 5.95, 2.5, 0.45, size=13, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ¿QUÉ ES LA PLATAFORMA?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, '¿Qué es RESPOND Guatemala?',
             'Plataforma web para registro prospectivo de trauma en hospitales de bajos recursos')

rect(sl, 0.3, 1.4, 5.8, 5.5, LIGHT_BLUE)
txt(sl, 'El problema', 0.5, 1.5, 5.4, 0.45, size=16, bold=True, color=NAVY)
problems = [
    'Los hospitales de Guatemala no tienen sistemas digitales de registro de trauma',
    'Los datos en papel no permiten análisis, benchmarking ni mejora de calidad',
    'Sin datos no hay evidencia para pedir recursos o mejorar protocolos',
    'Las soluciones comerciales cuestan $50,000–$200,000/año',
    'Los sistemas existentes no funcionan sin internet',
]
for i, p in enumerate(problems):
    txt(sl, f'▸  {p}', 0.45, 2.1 + i*0.82, 5.5, 0.75, size=11.5, color=DARK)

rect(sl, 6.4, 1.4, 6.6, 5.5, LIGHT_GREEN)
txt(sl, 'Nuestra solución', 6.6, 1.5, 6.2, 0.45, size=16, bold=True, color=GREEN)
solutions = [
    'App web gratuita, accesible desde cualquier dispositivo con navegador',
    'Funciona offline (PWA) — sincroniza cuando regresa el internet',
    'Formulario de 16 pasos optimizado para flujo clínico real',
    'Cálculo automático de GCS e ISS — sin tablas ni cálculos manuales',
    'Dashboard en tiempo real para el equipo de trauma',
    'Bilingüe español/inglés — lista para uso inmediato',
    'Costo de operación actual: $0/mes',
]
for i, s in enumerate(solutions):
    txt(sl, f'✓  {s}', 6.45, 2.1 + i*0.68, 6.3, 0.62, size=11.5, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ¿QUÉ CONTIENE?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, '¿Qué contiene la plataforma?', 'Módulos y funcionalidades del MVP actual')

features = [
    ('📋', 'Formulario\n16 pasos', LIGHT_BLUE),
    ('🧠', 'Cálculo GCS\nautomático', LIGHT_GREEN),
    ('📊', 'Score ISS\nautomático', LIGHT_BLUE),
    ('🗺️', 'Mapa corporal\ninteractivo', LIGHT_GREEN),
    ('📈', 'Dashboard\nen tiempo real', LIGHT_BLUE),
    ('📥', 'Exportar\na Excel', LIGHT_GREEN),
    ('📴', 'Modo offline\n(PWA)', LIGHT_BLUE),
    ('🌐', 'Bilingüe\nES / EN', LIGHT_GREEN),
    ('🔐', 'Auth con\nroles', LIGHT_BLUE),
    ('🏥', 'Multi-hospital\nlisto', LIGHT_GREEN),
    ('📱', 'Mobile\nfriendly', LIGHT_BLUE),
    ('🔄', 'Sincronización\nautomática', LIGHT_GREEN),
]

cols = 6
for i, (emoji, label, color) in enumerate(features):
    col = i % cols
    row = i // cols
    x = 0.3 + col * 2.15
    y = 1.45 + row * 1.9
    rect_border(sl, x, y, 1.98, 1.7, color, BLUE, 1)
    txt(sl, emoji, x, y+0.08, 1.98, 0.7, size=28, align=PP_ALIGN.CENTER)
    txt_box(sl, label, x, y+0.85, 1.98, 0.78, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CÓMO FUE CONSTRUIDA
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Arquitectura técnica actual',
             'Stack moderno, probado en producción, $0/mes de costo operativo')

# Left: Tech stack
rect(sl, 0.3, 1.4, 5.9, 5.5, LIGHT_GRAY)
txt(sl, 'Stack tecnológico', 0.5, 1.5, 5.5, 0.42, size=15, bold=True, color=NAVY)

stack = [
    ('⚡ Next.js 14', 'Framework frontend + backend (React)', BLUE),
    ('🔷 TypeScript', 'Tipado estático, menos bugs en producción', BLUE),
    ('🗄️ Supabase', 'Base de datos PostgreSQL + autenticación', BLUE),
    ('🎨 Tailwind CSS', 'Diseño responsive mobile-first', BLUE),
    ('🧩 Shadcn/UI', 'Componentes UI accesibles y consistentes', BLUE),
    ('☁️ Vercel', 'Hosting global, CDN, deploys automáticos', BLUE),
    ('📴 Service Worker', 'PWA offline-first con sincronización', BLUE),
]
for i, (tech, desc, c) in enumerate(stack):
    rect(sl, 0.35, 2.05+i*0.56, 5.8, 0.5, WHITE)
    txt(sl, tech, 0.45, 2.08+i*0.56, 2.1, 0.42, size=11.5, bold=True, color=NAVY)
    txt(sl, desc, 2.6, 2.1+i*0.56, 3.4, 0.38, size=10.5, color=DARK)

# Right: Architecture flow
rect(sl, 6.5, 1.4, 6.5, 5.5, LIGHT_BLUE)
txt(sl, 'Flujo de datos actual', 6.7, 1.5, 6.1, 0.42, size=15, bold=True, color=NAVY)

flow_items = [
    ('👤 Usuario (browser)', NAVY, WHITE),
    ('⬇', None, None),
    ('📱 App Next.js (Vercel)', BLUE, WHITE),
    ('⬇', None, None),
    ('🔐 Supabase Auth\n(JWT tokens)', RGBColor(0x1A, 0x5A, 0x9E), WHITE),
    ('⬇', None, None),
    ('🗄️ Supabase PostgreSQL\n(RLS policies)', RGBColor(0x0E, 0x6B, 0x5E), WHITE),
    ('⬇', None, None),
    ('📊 Dashboard / Excel Export', NAVY, WHITE),
]

y_pos = 2.0
for item, fill, tc in flow_items:
    if item == '⬇':
        txt(sl, '↓', 9.2, y_pos-0.02, 1.0, 0.32, size=18, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        y_pos += 0.32
    else:
        rect(sl, 6.7, y_pos, 6.0, 0.62, fill)
        txt_box(sl, item, 6.75, y_pos+0.05, 5.9, 0.52, size=11, bold=True, color=tc, align=PP_ALIGN.CENTER)
        y_pos += 0.72

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SEGURIDAD ACTUAL
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Seguridad de la plataforma',
             'Protecciones implementadas en el MVP actual')

# Already have
rect(sl, 0.3, 1.4, 6.1, 5.5, LIGHT_GREEN)
txt(sl, '✅  Seguridad implementada', 0.5, 1.5, 5.7, 0.45, size=15, bold=True, color=GREEN)

have = [
    ('🔒 HTTPS/TLS', 'Todo el tráfico cifrado en tránsito — certificado automático Vercel'),
    ('🔐 Autenticación JWT', 'Tokens firmados con rotación automática (Supabase Auth)'),
    ('🛡️ Row Level Security', 'Cada usuario solo ve sus propios pacientes (PostgreSQL RLS)'),
    ('💾 Cifrado en reposo', 'Base de datos cifrada AES-256 (Supabase/AWS)'),
    ('🌐 CDN Global', 'Vercel Edge Network — protección DDoS incluida'),
    ('🔑 Variables de entorno', 'Secretos en Vercel, nunca en el código fuente'),
    ('📦 SOC 2 Type II', 'Supabase certificado — auditorías anuales de seguridad'),
]
for i, (title, desc) in enumerate(have):
    rect(sl, 0.35, 2.05+i*0.65, 6.0, 0.6, WHITE)
    txt(sl, title, 0.45, 2.09+i*0.65, 2.1, 0.42, size=11, bold=True, color=GREEN)
    txt(sl, desc, 2.55, 2.12+i*0.65, 3.7, 0.38, size=10.5, color=DARK)

# Need for HIPAA
rect(sl, 6.7, 1.4, 6.3, 5.5, LIGHT_ORANGE)
txt(sl, '⚠️  Pendiente para HIPAA', 6.9, 1.5, 5.9, 0.45, size=15, bold=True, color=ORANGE)

need = [
    ('⏱️ Session timeout', 'Auto-logout tras 15 min de inactividad'),
    ('📋 Audit log', 'Registro de quién accedió/modificó qué paciente y cuándo'),
    ('📱 MFA', 'Autenticación de 2 factores para todos los usuarios'),
    ('🔏 BAA firmado', 'Contrato legal con el proveedor de base de datos'),
    ('📄 Políticas de privacidad', 'Aviso de prácticas de privacidad publicado'),
    ('📊 Risk Assessment', 'Documentación formal de evaluación de riesgos'),
    ('🎓 Training records', 'Registro de entrenamiento del personal'),
]
for i, (title, desc) in enumerate(need):
    rect(sl, 6.75, 2.05+i*0.65, 6.2, 0.6, WHITE)
    txt(sl, title, 6.85, 2.09+i*0.65, 2.1, 0.42, size=11, bold=True, color=ORANGE)
    txt(sl, desc, 9.0, 2.12+i*0.65, 3.8, 0.38, size=10.5, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ¿POR QUÉ HIPAA?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, '¿Por qué necesitamos cumplir con HIPAA?',
             'Requerimiento de los mentores para continuar el proyecto')

rect(sl, 0.3, 1.4, 4.1, 5.5, LIGHT_RED)
txt(sl, '⚖️  ¿Qué es HIPAA?', 0.5, 1.5, 3.7, 0.45, size=14, bold=True, color=RED)
txt_box(sl,
    'Health Insurance Portability\nand Accountability Act\n(Ley federal de EE.UU., 1996)\n\n'
    'Protege la información médica\nidentificable de pacientes\n(PHI — Protected Health Information)\n\n'
    'Requiere:\n• Confidencialidad\n• Integridad\n• Disponibilidad\nde todos los datos del paciente',
    0.45, 2.05, 3.7, 4.6, size=11.5, color=DARK)

rect(sl, 4.65, 1.4, 4.0, 5.5, LIGHT_ORANGE)
txt(sl, '🌎  ¿Aplica en Guatemala?', 4.85, 1.5, 3.6, 0.45, size=13, bold=True, color=ORANGE)
txt_box(sl,
    'Técnicamente NO — HIPAA\nes ley estadounidense.\n\n'
    'PERO aplica cuando:\n\n'
    '• Los mentores son de una\n  institución de EE.UU. (BU)\n\n'
    '• El financiamiento es de\n  NIH u otras fuentes de EE.UU.\n\n'
    '• Se planea publicar datos\n  en revistas que requieren\n  estándares internacionales\n\n'
    '→ Cumplir con HIPAA es\n  el estándar de oro para\n  investigación clínica seria',
    4.85, 2.05, 3.7, 4.6, size=11.5, color=DARK)

rect(sl, 8.9, 1.4, 4.1, 5.5, LIGHT_BLUE)
txt(sl, '📋  ¿Qué protege?', 9.1, 1.5, 3.7, 0.45, size=14, bold=True, color=NAVY)
phi_items = [
    'Nombre del paciente',
    'Fecha de nacimiento',
    'Número de expediente',
    'Fecha de admisión/alta',
    'Diagnóstico',
    'Resultados clínicos (GCS, ISS)',
    'Procedimientos realizados',
    'Condición al egreso',
    'Cualquier dato que permita\nidentificar al paciente',
]
for i, item in enumerate(phi_items):
    txt(sl, f'▸  {item}', 9.1, 2.08+i*0.52, 3.7, 0.48, size=11, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — OPCIÓN 1: REDCAP (RECOMENDADA)
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Opción 1: Integración con REDCap  ⭐ RECOMENDADA',
             'La app en Vercel envía datos directamente al servidor REDCap institucional')

# Header badge
rect(sl, 10.8, 1.35, 2.2, 0.55, GREEN)
txt(sl, '⭐ RECOMENDADA', 10.82, 1.38, 2.16, 0.48, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Architecture flow (horizontal)
rect(sl, 0.3, 1.45, 12.73, 2.3, LIGHT_GRAY)
txt(sl, 'Flujo de datos con REDCap', 0.5, 1.5, 8.0, 0.38, size=13, bold=True, color=NAVY)

boxes = [
    ('👤\nEnfermera\n/ Médico', NAVY, WHITE, 0.4, 1.95, 1.7, 1.6),
    ('📱\nApp RESPOND\n(Vercel)', BLUE, WHITE, 2.45, 1.95, 1.9, 1.6),
    ('🔀\nAPI Route\nNext.js', RGBColor(0x5D,0x6D,0x7E), WHITE, 4.65, 1.95, 1.8, 1.6),
    ('🏛️\nREDCap API\n(institucional)', TEAL, WHITE, 6.75, 1.95, 2.0, 1.6),
    ('🗄️\nServidor REDCap\nBU / Institución', RGBColor(0x0E,0x5E,0x5E), WHITE, 9.05, 1.95, 2.1, 1.6),
    ('🔒\nHIPAA\nCubierto', GREEN, WHITE, 11.4, 1.95, 1.6, 1.6),
]
for emoji_label, fill, tc, x, y, w, h in boxes:
    rect(sl, x, y, w, h, fill)
    txt_box(sl, emoji_label, x, y+0.15, w, h-0.2, size=10.5, bold=True, color=tc, align=PP_ALIGN.CENTER)

for ax in [2.2, 4.5, 6.6, 8.9, 11.25]:
    txt(sl, '→', ax, 2.5, 0.35, 0.42, size=18, bold=True, color=BLUE)

# Two columns below
rect(sl, 0.3, 3.95, 6.1, 3.25, LIGHT_GREEN)
txt(sl, '✅  Ventajas', 0.5, 4.05, 5.7, 0.4, size=14, bold=True, color=GREEN)
pros = [
    ('💰 Costo $0/año', 'Si la institución ya tiene REDCap (BU lo tiene)'),
    ('🔒 HIPAA cubierto', 'La institución es responsable legal, no tú'),
    ('📊 Exportaciones', 'SPSS, Stata, R, CSV, Excel — nativas en REDCap'),
    ('🎓 Estándar en research', '6,000+ instituciones lo usan mundialmente'),
    ('🚀 Sin BAA que firmar', 'El departamento de IT institucional lo gestiona'),
]
for i, (title, desc) in enumerate(pros):
    txt(sl, title, 0.45, 4.58+i*0.52, 1.95, 0.45, size=11, bold=True, color=GREEN)
    txt(sl, desc, 2.45, 4.6+i*0.52, 3.8, 0.42, size=10.5, color=DARK)

rect(sl, 6.7, 3.95, 6.3, 3.25, LIGHT_ORANGE)
txt(sl, '⚠️  Consideraciones', 6.9, 4.05, 5.9, 0.4, size=14, bold=True, color=ORANGE)
cons = [
    ('🏛️ Requiere institución', 'Necesitas que BU u otra institución abra el proyecto'),
    ('🗺️ Mapeo de campos', 'Los campos del form deben mapearse a variables REDCap'),
    ('📴 Offline limitado', 'REDCap API requiere conexión (soluble con caché local)'),
    ('📈 Dashboard externo', 'REDCap reports son básicos — tu app puede pullar los datos'),
    ('⏱️ Setup inicial', '~1 semana para coordinar con IT + mapear variables'),
]
for i, (title, desc) in enumerate(cons):
    txt(sl, title, 6.9, 4.58+i*0.52, 2.1, 0.45, size=11, bold=True, color=ORANGE)
    txt(sl, desc, 9.0, 4.6+i*0.52, 3.85, 0.42, size=10.5, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — REDCAP CÓMO FUNCIONA TÉCNICAMENTE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'REDCap — Integración técnica',
             'La app ya existente se conecta al API de REDCap con cambios mínimos')

rect(sl, 0.3, 1.4, 5.9, 5.5, LIGHT_BLUE)
txt(sl, '🔧 ¿Qué cambia en el código?', 0.5, 1.5, 5.5, 0.42, size=14, bold=True, color=NAVY)

changes = [
    ('Muy poco', 'El formulario de 16 pasos no cambia'),
    ('Muy poco', 'El dashboard puede seguir leyendo de REDCap API'),
    ('Nuevo', 'API route en Next.js que recibe el form y hace POST a REDCap'),
    ('Nuevo', 'Mapeo de campos: nombre_campo_app → redcap_variable_name'),
    ('Nuevo', 'Variable de entorno: REDCAP_API_TOKEN + REDCAP_URL'),
    ('Opcional', 'Caché local (IndexedDB) para modo offline → sync al reconectar'),
]
labels = {'Muy poco': GREEN, 'Nuevo': BLUE, 'Opcional': ORANGE}
for i, (label, desc) in enumerate(changes):
    lc = labels.get(label, GRAY)
    rect(sl, 0.35, 2.1+i*0.75, 1.2, 0.55, lc)
    txt(sl, label, 0.37, 2.12+i*0.75, 1.16, 0.5, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, desc, 1.65, 2.15+i*0.75, 4.4, 0.5, size=11, color=DARK)

rect(sl, 6.5, 1.4, 6.5, 5.5, LIGHT_GRAY)
txt(sl, '📋 Proceso de setup', 6.7, 1.5, 6.1, 0.42, size=14, bold=True, color=NAVY)

steps = [
    ('1', 'Solicitar proyecto REDCap en BU (IT dept.)', '~1-3 días'),
    ('2', 'Definir variables REDCap = campos del form', '~1 día'),
    ('3', 'Generar API token en REDCap dashboard', '5 min'),
    ('4', 'Escribir la API route en Next.js', '~4-8 horas'),
    ('5', 'Mapear 50-60 campos del form al API', '~4-8 horas'),
    ('6', 'Testing end-to-end (form → REDCap)', '~1 día'),
    ('7', 'Deploy + capacitación del equipo', '~1 día'),
]
for i, (num, step, time) in enumerate(steps):
    rect(sl, 6.55, 2.1+i*0.65, 0.45, 0.52, NAVY)
    txt(sl, num, 6.57, 2.13+i*0.65, 0.41, 0.45, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, step, 7.1, 2.15+i*0.65, 4.2, 0.42, size=11, color=DARK)
    txt(sl, time, 11.35, 2.15+i*0.65, 1.5, 0.42, size=10, color=GRAY, italic=True, align=PP_ALIGN.RIGHT)

rect(sl, 6.55, 6.6-0.2, 6.4, 0.48, GREEN)
txt(sl, '⏱️  Tiempo total estimado: 1–2 semanas  ·  Costo de desarrollo: mínimo',
    6.65, 6.63-0.2, 6.2, 0.38, size=11.5, bold=True, color=WHITE)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — OPCIÓN 2: SUPABASE PRO
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Opción 2: Supabase Pro + BAA',
             'Upgrade del backend actual — mínimo cambio de código, HIPAA compliant')

cost_badge(sl, '$300', 'por año', 0.3, 1.4, BLUE)
rect(sl, 2.7, 1.4, 10.3, 0.9, LIGHT_BLUE)
txt_box(sl,
    'La app ya está 100% construida sobre Supabase. Solo se hace upgrade al plan Pro ($25/mes) y se firma el BAA (Business Associate Agreement). No hay migración de código.',
    2.85, 1.45, 10.0, 0.8, size=13, color=DARK)

# Architecture
rect(sl, 0.3, 2.45, 12.73, 1.5, LIGHT_GRAY)
txt(sl, 'Arquitectura (sin cambios):', 0.5, 2.52, 4.0, 0.38, size=12, bold=True, color=NAVY)
flow_items2 = [
    ('👤 Usuario', NAVY), ('📱 App Vercel', BLUE), ('🔐 Supabase Auth', TEAL),
    ('🗄️ Supabase PostgreSQL', RGBColor(0x0E,0x5E,0x5E)), ('✅ BAA firmado = HIPAA', GREEN)
]
for i, (label, fill) in enumerate(flow_items2):
    rect(sl, 0.3+i*2.55, 2.95, 2.4, 0.7, fill)
    txt(sl, label, 0.33+i*2.55, 2.98, 2.34, 0.64, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.78+i*2.55, 3.17, 0.35, 0.35, size=16, bold=True, color=BLUE)

cols_y = 4.1
rect(sl, 0.3, cols_y, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Ventajas', 0.5, cols_y+0.1, 5.7, 0.4, size=14, bold=True, color=GREEN)
pros2 = [
    'Cero cambios en el código de la app',
    'BAA disponible en plan Pro — se firma en el dashboard',
    'Supabase corre sobre AWS us-east-1 (misma infraestructura que hospitales grandes)',
    'Auth, DB, RLS, storage — todo HIPAA en un solo proveedor',
    'Dashboard, offline PWA y exportaciones funcionan igual',
]
for i, p in enumerate(pros2):
    txt(sl, f'✓  {p}', 0.45, cols_y+0.65+i*0.48, 5.7, 0.43, size=11.5, color=DARK)

rect(sl, 6.7, cols_y, 6.3, 3.05, LIGHT_ORANGE)
txt(sl, '⚠️  Consideraciones', 6.9, cols_y+0.1, 5.9, 0.4, size=14, bold=True, color=ORANGE)
cons2 = [
    'Costo fijo $25/mes aunque tengas 0 pacientes ese mes',
    'Vercel sigue siendo gratuito (PHI no pasa por Vercel)',
    'Requiere configurar session timeout + audit log + MFA (~1-2 días)',
    'Documentación administrativa aún necesaria (risk assessment, policies)',
]
for i, c in enumerate(cons2):
    txt(sl, f'▸  {c}', 6.9, cols_y+0.65+i*0.62, 5.9, 0.55, size=11.5, color=DARK)

# Steps
rect(sl, 0.3, 6.75-0.55, 12.73, 0.42, NAVY)
steps_txt = '  PASOS:   1. Upgrade Supabase → Pro   →   2. Firmar BAA en dashboard   →   3. Configurar session timeout + MFA + audit log   →   4. Documentar risk assessment'
txt(sl, steps_txt, 0.4, 6.78-0.55, 12.5, 0.35, size=11, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — OPCIÓN 3: AWS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Opción 3: Amazon Web Services (AWS)',
             'Infraestructura de clase empresarial — máxima flexibilidad y control')

cost_badge(sl, '$600–$1,800', 'por año (variable)', 0.3, 1.4, RGBColor(0xFF,0x99,0x00))
rect(sl, 3.0, 1.4, 10.0, 0.9, LIGHT_ORANGE)
txt_box(sl,
    'AWS firma BAA en cualquier plan pago. RDS PostgreSQL + Cognito (auth) + Amplify (hosting).\nMigración significativa de código. Usado por Epic, Cerner y sistemas hospitalarios grandes.',
    3.15, 1.45, 9.7, 0.8, size=12.5, color=DARK)

# Architecture
rect(sl, 0.3, 2.45, 12.73, 1.5, LIGHT_GRAY)
txt(sl, 'Arquitectura AWS:', 0.5, 2.52, 4.0, 0.38, size=12, bold=True, color=NAVY)
aws_items = [
    ('👤 Usuario', NAVY), ('📱 AWS Amplify\n(hosting)', ORANGE),
    ('🔐 Cognito\n(auth)', RGBColor(0xC0,0x50,0x00)),
    ('🗄️ RDS PostgreSQL\n+ KMS encrypt', RGBColor(0x8B,0x35,0x00)),
    ('✅ BAA = HIPAA', GREEN)
]
for i, (label, fill) in enumerate(aws_items):
    rect(sl, 0.3+i*2.55, 2.9, 2.4, 0.82, fill)
    txt_box(sl, label, 0.33+i*2.55, 2.93, 2.34, 0.76, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.78+i*2.55, 3.17, 0.35, 0.35, size=16, bold=True, color=ORANGE)

rect(sl, 0.3, 4.1, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Ventajas', 0.5, 4.2, 5.7, 0.4, size=14, bold=True, color=GREEN)
pros3 = [
    'Estándar de oro para healthcare en EE.UU.',
    'BAA disponible — AWS es HIPAA-eligible en todos los servicios clave',
    'Máximo control sobre la infraestructura',
    'Escala a millones de pacientes sin cambiar arquitectura',
    'Auditoría, logging (CloudTrail) y compliance nativos',
]
for i, p in enumerate(pros3):
    txt(sl, f'✓  {p}', 0.45, 4.75+i*0.48, 5.7, 0.43, size=11.5, color=DARK)

rect(sl, 6.7, 4.1, 6.3, 3.05, LIGHT_RED)
txt(sl, '✗  Desventajas', 6.9, 4.2, 5.9, 0.4, size=14, bold=True, color=RED)
cons3 = [
    'Migración de todo el código — auth, DB, API — semanas de trabajo',
    'Complejidad operativa alta: IAM, VPC, security groups',
    'Requiere experiencia DevOps para mantener',
    'Costo variable — puede subir con el uso',
    'Overkill para un proyecto de un hospital en etapa piloto',
]
for i, c in enumerate(cons3):
    txt(sl, f'✗  {c}', 6.9, 4.75+i*0.48, 5.9, 0.43, size=11.5, color=DARK)

# Cost breakdown
rect(sl, 0.3, 6.75-0.55, 12.73, 0.42, RGBColor(0xCC,0x77,0x00))
txt(sl, '  Costo estimado: RDS db.t3.micro ~$15/mes · Cognito <1K usuarios gratis · Amplify ~$1-5/mes · Extras (backups, KMS, CloudWatch) ~$5-10/mes → Total ~$25-35/mes',
    0.4, 6.78-0.55, 12.5, 0.35, size=10.5, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — OPCIÓN 4: GOOGLE CLOUD / FIREBASE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Opción 4: Google Cloud Platform / Firebase',
             'Flexible y de bajo costo a volúmenes bajos — BAA disponible en plan Blaze')

cost_badge(sl, '$0–$360', 'por año (pay-as-you-go)', 0.3, 1.4, TEAL)
rect(sl, 3.0, 1.4, 10.0, 0.9, LIGHT_TEAL)
txt_box(sl,
    'Google firma BAA bajo Google Cloud Platform. Firebase Auth + Firestore (o Cloud SQL) son HIPAA-eligible.\nCosto casi $0 a volúmenes bajos. Migración significativa de código (NoSQL o nuevo ORM).',
    3.15, 1.45, 9.7, 0.8, size=12.5, color=DARK)

# Architecture
rect(sl, 0.3, 2.45, 12.73, 1.5, LIGHT_GRAY)
txt(sl, 'Arquitectura GCP:', 0.5, 2.52, 4.0, 0.38, size=12, bold=True, color=NAVY)
gcp_items = [
    ('👤 Usuario', NAVY),
    ('📱 Cloud Run\n(hosting)', TEAL),
    ('🔐 Firebase Auth\n(Google)', RGBColor(0x0A,0x70,0x70)),
    ('🗄️ Firestore o\nCloud SQL', RGBColor(0x06,0x50,0x50)),
    ('✅ BAA = HIPAA', GREEN)
]
for i, (label, fill) in enumerate(gcp_items):
    rect(sl, 0.3+i*2.55, 2.9, 2.4, 0.82, fill)
    txt_box(sl, label, 0.33+i*2.55, 2.93, 2.34, 0.76, size=10.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, '→', 2.78+i*2.55, 3.17, 0.35, 0.35, size=16, bold=True, color=TEAL)

rect(sl, 0.3, 4.1, 6.1, 3.05, LIGHT_GREEN)
txt(sl, '✅  Ventajas', 0.5, 4.2, 5.7, 0.4, size=14, bold=True, color=GREEN)
pros4 = [
    'Costo casi $0 con volúmenes bajos de datos',
    'Google firma BAA en plan Blaze (gratuito o pay-as-you-go)',
    'Firebase Auth muy sólido con MFA incluido',
    'Cloud SQL = PostgreSQL (más fácil migrar que Firestore)',
    'Integración con Google Workspace / Drive para reportes',
]
for i, p in enumerate(pros4):
    txt(sl, f'✓  {p}', 0.45, 4.75+i*0.48, 5.7, 0.43, size=11.5, color=DARK)

rect(sl, 6.7, 4.1, 6.3, 3.05, LIGHT_RED)
txt(sl, '✗  Desventajas', 6.9, 4.2, 5.9, 0.4, size=14, bold=True, color=RED)
cons4 = [
    'Migración completa del código (Supabase → Firebase/GCP)',
    'Firestore es NoSQL — estructura muy distinta al esquema actual',
    'Cloud SQL requiere configurar más piezas que Supabase',
    'Curva de aprendizaje de GCP vs. la simplicidad de Supabase',
    'Vendor lock-in en ecosistema Google',
]
for i, c in enumerate(cons4):
    txt(sl, f'✗  {c}', 6.9, 4.75+i*0.48, 5.9, 0.43, size=11.5, color=DARK)

rect(sl, 0.3, 6.75-0.55, 12.73, 0.42, TEAL)
txt(sl, '  Costo estimado: Firestore <50K lecturas/día gratis · Firebase Auth gratis · Cloud Run $0-5/mes a bajo volumen → Total ~$0-30/mes',
    0.4, 6.78-0.55, 12.5, 0.35, size=10.5, color=WHITE, bold=True)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPARACIÓN DE COSTOS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Comparación de costos y características',
             'Análisis anual para un hospital piloto con ~500-1,000 pacientes/año')

# Cost badges row
options = [
    ('REDCap\n(Institucional)', '$0', '/año', GREEN),
    ('Supabase Pro', '$300', '/año', BLUE),
    ('Google Cloud\n(Firebase)', '$0–360', '/año', TEAL),
    ('AWS\n(RDS + Cognito)', '$300–600', '/año', ORANGE),
]
for i, (name, cost, unit, color) in enumerate(options):
    x = 0.3 + i * 3.25
    rect(sl, x, 1.4, 3.0, 1.1, color)
    txt_box(sl, name, x, 1.42, 3.0, 0.42, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, cost, x, 1.82, 3.0, 0.55, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, unit, x, 2.37, 3.0, 0.28, size=11, color=RGBColor(0xDD,0xDD,0xDD), align=PP_ALIGN.CENTER)

# Comparison table
headers = ['Característica', 'REDCap', 'Supabase Pro', 'Google Cloud', 'AWS']
col_x = [0.3, 3.55, 6.05, 8.55, 11.05]
col_w = [3.2, 2.45, 2.45, 2.45, 2.2]

# Header row
rect(sl, 0.3, 2.7, 12.73, 0.48, NAVY)
for j, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    txt(sl, h, x+0.05, 2.73, w-0.1, 0.4, size=12, bold=True, color=WHITE,
        align=PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT)

rows = [
    ('Cumplimiento HIPAA',         '✅ Institucional',  '✅ BAA Pro',      '✅ BAA Blaze',    '✅ BAA cualquier plan'),
    ('Costo anual estimado',       '💚 $0',             '🔵 $300',         '🟢 $0–$360',      '🟡 $300–$600'),
    ('Cambio de código requerido', '🟡 Medio (API)',    '✅ Mínimo',       '🔴 Alto',         '🔴 Alto'),
    ('Tiempo de implementación',   '1–2 semanas',       '2–3 días',        '4–8 semanas',     '4–8 semanas'),
    ('Offline / PWA',              '🟡 Con caché local','✅ Completo',     '🟡 Con caché',    '🟡 Con caché'),
    ('Dashboard en tiempo real',   '🟡 Con API pull',   '✅ Ya funciona',  '🟡 Requiere dev', '🟡 Requiere dev'),
    ('Exportar para análisis',     '✅ SPSS/SAS/R/CSV', '✅ Excel/CSV',    '🟡 CSV/JSON',     '🟡 CSV/JSON'),
    ('Conocido por investigadores','✅ Estándar de oro','🟡 No conocido',  '🟡 No conocido',  '🟡 No conocido'),
    ('Escalabilidad',              '✅ Alta',           '✅ Alta',         '✅ Alta',          '✅ Muy alta'),
]

row_colors = [LIGHT_GRAY, WHITE]
for i, row in enumerate(rows):
    bg_c = row_colors[i % 2]
    rect(sl, 0.3, 3.22+i*0.42, 12.73, 0.41, bg_c)
    for j, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        al = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        txt(sl, cell, x+0.05, 3.24+i*0.42, w-0.08, 0.36, size=10, color=DARK, align=al)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — RECOMENDACIÓN
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Recomendación',
             'Estrategia en dos fases según el objetivo del proyecto')

rect(sl, 0.3, 1.4, 6.1, 5.55, LIGHT_GREEN)
rect(sl, 0.3, 1.4, 6.1, 0.62, GREEN)
txt(sl, '⭐  Fase 1 — Piloto de investigación', 0.45, 1.44, 5.8, 0.52, size=15, bold=True, color=WHITE)

txt(sl, 'REDCap + App RESPOND Guatemala', 0.45, 2.15, 5.8, 0.48, size=16, bold=True, color=NAVY)
txt(sl, '$0/año si BU tiene REDCap', 0.45, 2.63, 5.8, 0.38, size=13, color=GREEN, bold=True)

rec1_items = [
    '✓  HIPAA cubierto por la institución sin costo',
    '✓  Datos en servidor de BU — tranquilidad legal total',
    '✓  REDCap es el estándar en publicaciones clínicas',
    '✓  Exportación directa a SPSS/Stata/R para análisis',
    '✓  IRB de BU puede aprobar directamente',
    '✓  Mínimo costo de operación mensual',
    '✓  1–2 semanas de integración técnica',
]
for i, item in enumerate(rec1_items):
    txt(sl, item, 0.45, 3.15+i*0.52, 5.7, 0.46, size=12, color=DARK)

rect(sl, 6.7, 1.4, 6.3, 5.55, LIGHT_BLUE)
rect(sl, 6.7, 1.4, 6.3, 0.62, BLUE)
txt(sl, '🚀  Fase 2 — Escalamiento regional', 6.85, 1.44, 6.0, 0.52, size=15, bold=True, color=WHITE)

txt(sl, 'Supabase Pro (o AWS si hay funding)', 6.85, 2.15, 6.0, 0.48, size=15, bold=True, color=NAVY)
txt(sl, '$300/año (Supabase) · $300–600/año (AWS)', 6.85, 2.63, 6.0, 0.38, size=12, color=BLUE, bold=True)

rec2_items = [
    '✓  Cuando haya múltiples hospitales (>3)',
    '✓  Si se necesita dashboard en tiempo real avanzado',
    '✓  Si el proyecto obtiene financiamiento formal',
    '✓  App ya construida — solo upgrade + BAA',
    '✓  Independiente de la institución académica',
    '✓  Plataforma propia, control total de los datos',
    '✓  2–3 días para implementar desde hoy',
]
for i, item in enumerate(rec2_items):
    txt(sl, item, 6.85, 3.15+i*0.52, 6.1, 0.46, size=12, color=DARK)

footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PRÓXIMOS PASOS
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, WHITE)
slide_header(sl, 'Próximos pasos', 'Camino hacia el lanzamiento del piloto con cumplimiento HIPAA')

steps_main = [
    (NAVY,  '1', 'Confirmar acceso a REDCap',
     'Preguntar a BU IT si existe proyecto REDCap disponible para Guatemala.\nAlternativa: solicitar nuevo proyecto (formulario + aprobación ~1-3 días).'),
    (BLUE,  '2', 'IRB / Comité de ética',
     'Someter protocolo de investigación al IRB de BU.\nLa integración con REDCap facilita la aprobación (es la plataforma estándar de BU).'),
    (TEAL,  '3', 'Integración técnica REDCap',
     'Mapear los 16 pasos del formulario a variables REDCap.\nEscribir API route en Next.js para enviar datos (4-8 horas de desarrollo).'),
    (GREEN, '4', 'Agregar controles HIPAA técnicos',
     'Session timeout (15 min) · Audit log · MFA para usuarios.\nEstimado: 1-2 días de desarrollo, independiente del backend.'),
    (ORANGE,'5', 'Hospital piloto + capacitación',
     'Agregar 1 hospital al dropdown · Crear 3-5 cuentas de usuario.\nCapacitación del equipo clínico: ~2 horas.'),
]

for i, (color, num, title, desc) in enumerate(steps_main):
    y = 1.45 + i * 1.12
    rect(sl, 0.3, y, 0.55, 0.95, color)
    txt(sl, num, 0.3, y+0.12, 0.55, 0.72, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(sl, 0.9, y, 11.9, 0.95, LIGHT_GRAY if i%2==0 else WHITE)
    txt(sl, title, 1.0, y+0.05, 5.0, 0.42, size=14, bold=True, color=color)
    txt_box(sl, desc, 1.0, y+0.48, 11.6, 0.44, size=11, color=DARK)

rect(sl, 0.3, 7.05, 12.73, 0.15, NAVY)
footer(sl)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CIERRE
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl, NAVY)
rect(sl, 0, 0, 0.18, 7.5, BLUE)
rect(sl, 0, 5.9, 13.33, 0.08, BLUE)

txt(sl, 'RESPOND Guatemala', 0.5, 0.8, 12.0, 0.75, size=36, bold=True, color=WHITE)
txt(sl, 'Registro Digital de Trauma · Listo para el piloto', 0.5, 1.6, 12.0, 0.55, size=20,
    color=RGBColor(0xAE,0xD6,0xF1), italic=True)

rect(sl, 0.5, 2.4, 12.0, 1.85, RGBColor(0x1E,0x49,0x8A))
summary = [
    ('✅', 'Plataforma MVP completa y en producción — trauma-registry.vercel.app'),
    ('✅', 'Seguridad base implementada: HTTPS, JWT, RLS, cifrado, SOC 2'),
    ('✅', 'Opciones HIPAA identificadas: REDCap ($0) · Supabase Pro ($300/año) · AWS · GCP'),
    ('⭐', 'Recomendación: REDCap institucional para piloto → mínimo costo, máxima credibilidad'),
    ('🚀', 'Tiempo al primer paciente registrado: 2–4 semanas con IRB aprobado'),
]
for i, (icon, text) in enumerate(summary):
    txt(sl, icon, 0.7, 2.52+i*0.33, 0.45, 0.3, size=13, bold=True, color=YELLOW)
    txt(sl, text, 1.2, 2.54+i*0.33, 11.0, 0.3, size=12, color=WHITE)

txt(sl, '¿Preguntas?', 0.5, 4.5, 12.0, 0.65, size=28, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
txt(sl, 'trauma-registry.vercel.app/en/patients', 0.5, 5.2, 12.0, 0.45,
    size=15, color=RGBColor(0x5D,0x8A,0xC4), italic=True, align=PP_ALIGN.CENTER)

rect(sl, 0, 6.05, 13.33, 1.45, RGBColor(0x12,0x27,0x4A))
txt(sl, 'RESPOND Guatemala  ·  Registro de Trauma  ·  Confidencial  ·  Marzo 2026',
    0.5, 6.55, 12.33, 0.45, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Total slides: {len(prs.slides)}')
